"""
YuvaSetu Presentation Generator
Generates a 10-slide professional PowerPoint for Hack DTU national hackathon.
Light background, dark fonts, architecture diagrams, statistics visualizations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# ── Color Palette (light bg, dark text) ──────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_GRAY  = RGBColor(0xE9, 0xEC, 0xEF)
MED_GRAY    = RGBColor(0xAD, 0xB5, 0xBD)
DARK_GRAY   = RGBColor(0x49, 0x50, 0x57)
NEAR_BLACK  = RGBColor(0x21, 0x25, 0x29)
BLACK       = RGBColor(0x00, 0x00, 0x00)

PRIMARY     = RGBColor(0x22, 0x8B, 0xE6)  # Blue
SECONDARY   = RGBColor(0x6C, 0x75, 0x7D)
SUCCESS     = RGBColor(0x2D, 0xCE, 0x89)  # Green
DANGER      = RGBColor(0xF5, 0x36, 0x5C)  # Red
WARNING     = RGBColor(0xFB, 0x6D, 0x40)  # Orange
ACCENT      = RGBColor(0x11, 0xCE, 0xF0)  # Cyan
PURPLE      = RGBColor(0x6F, 0x42, 0xC1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=NEAR_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, left, top, width, height, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def set_text(tf, text, size=14, bold=False, color=NEAR_BLACK, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_para(tf, text, size=14, bold=False, color=NEAR_BLACK, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2), font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_title_bar(slide, title, subtitle=""):
    """Top accent bar with title."""
    bar = add_shape_box(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), PRIMARY)
    tf = bar.text_frame
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.1)
    tf.word_wrap = True
    set_text(tf, title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_para(tf, subtitle, size=14, color=RGBColor(0xCC, 0xDD, 0xFF), space_before=Pt(0))


def add_stat_card(slide, left, top, number, label, color=PRIMARY, width=Inches(2.4), height=Inches(1.5)):
    card = add_shape_box(slide, left, top, width, height, WHITE, border_color=LIGHT_GRAY, border_width=Pt(1))
    tf = card.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True
    set_text(tf, number, size=32, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_para(tf, label, size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(2))
    return card


def add_flow_box(slide, left, top, title, subtitle, color=PRIMARY, w=Inches(1.8), h=Inches(1.1)):
    box = add_shape_box(slide, left, top, w, h, color)
    tf = box.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.08)
    tf.word_wrap = True
    set_text(tf, title, size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_para(tf, subtitle, size=9, color=RGBColor(0xDD, 0xEE, 0xFF), alignment=PP_ALIGN.CENTER, space_before=Pt(2))
    return box


def add_table(slide, left, top, width, height, rows_data, col_widths=None, header_color=PRIMARY):
    """Add a styled table. rows_data = list of lists. First row = header."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = "Segoe UI"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.size = Pt(10)
                else:
                    paragraph.font.color.rgb = NEAR_BLACK

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF3, 0xF5)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


def add_bar_chart(slide, left, top, width, height, categories, values, title="", bar_color=PRIMARY):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = NEAR_BLACK

    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = bar_color

    return chart_frame


# ═══════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════

def slide_01_problem(prs):
    """SLIDE 1: Problem + SDG Mapping"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "The Problem — India's Employment Crisis", "PS 01: Next-Gen Generative AI  |  Hack DTU")

    # Stat cards
    stats = [
        ("1.5M", "Engg graduates/year\n(only 10% get jobs)", DANGER),
        ("83%", "Still unemployed\nafter graduating", WARNING),
        ("82%", "Employers report\nskill mismatch", PRIMARY),
        ("88%", "ATS rejects qualified\ntalent (HBS 2021)", PURPLE),
    ]
    for i, (num, lbl, clr) in enumerate(stats):
        add_stat_card(slide, Inches(0.5 + i * 3.1), Inches(1.25), num, lbl, clr, width=Inches(2.8), height=Inches(1.4))

    # 4 Gaps
    gap_box = add_shape_box(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.4), PRIMARY)
    set_text(gap_box.text_frame, "  4 Gaps No Platform Solves Together", size=14, bold=True, color=WHITE)

    gaps = [
        ("ACCESS", "Coaching costs Rs 5K-50K/session.\n80% students can't afford it.", DANGER),
        ("FEEDBACK", "Students fail interviews but\nnever learn WHY.", WARNING),
        ("MATCHING", "Resumes sent blindly.\nKeyword ATS misses talent.", PRIMARY),
        ("LANGUAGE", "English-only platforms.\n60% of India excluded.", PURPLE),
    ]
    for i, (title, desc, clr) in enumerate(gaps):
        bx = add_shape_box(slide, Inches(0.5 + i * 3.1), Inches(3.45), Inches(2.8), Inches(1.15), WHITE, border_color=clr, border_width=Pt(2))
        tf = bx.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.word_wrap = True
        set_text(tf, title, size=13, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, desc, size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))

    # SDG Mapping
    sdg_header = add_shape_box(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4), RGBColor(0x1B, 0x4D, 0x3E))
    set_text(sdg_header.text_frame, "  SDG Alignment + Problem Statement Mapping", size=14, bold=True, color=WHITE)

    sdg_data = [
        ["SDG Goal", "PS Mandate", "Our Feature", "Impact"],
        ["SDG 4: Quality Education", "Democratize education", "AI Mock Interview (Rs 0) + 237 free govt courses", "2000x cheaper"],
        ["SDG 8: Decent Work", "Automate workflows", "Semantic matching (<100ms) + auto-scored feedback", "88% hidden talent found"],
        ["SDG 10: Reduced Inequalities", "Bridge language barriers", "Voice interviews in Hindi/Marathi/English", "60% more India covered"],
    ]
    add_table(slide, Inches(0.5), Inches(5.35), Inches(12.3), Inches(1.8), sdg_data,
              col_widths=[Inches(2.5), Inches(2.5), Inches(4.5), Inches(2.8)],
              header_color=RGBColor(0x1B, 0x4D, 0x3E))


def slide_02_tech(prs):
    """SLIDE 2: Tech Stack + 6 AI Systems"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Tech Architecture — 6 AI Systems", "Full-stack: React + Express + MongoDB + 3 AI microservices")

    # -- Architecture diagram --
    # Frontend box
    fe = add_shape_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(1.5), RGBColor(0xDB, 0xEA, 0xFE), border_color=PRIMARY, border_width=Pt(2))
    tf = fe.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, "FRONTEND", size=14, bold=True, color=PRIMARY)
    add_para(tf, "React 18  ·  TypeScript  ·  Vite  ·  TailwindCSS", size=10, color=DARK_GRAY)
    add_para(tf, "shadcn/ui  ·  i18next (EN/HI/MR)  ·  Socket.IO", size=10, color=DARK_GRAY)
    add_para(tf, "face-api.js (proctoring)  ·  pdfjs-dist (resume)", size=10, color=DARK_GRAY)
    add_para(tf, "28,680 LOC  ·  124 components  ·  1,162 translation keys", size=10, bold=True, color=PRIMARY)

    # Arrow
    add_arrow(slide, Inches(2.5), Inches(2.85), Inches(0.8), Inches(0.3), PRIMARY)

    # Backend box
    be = add_shape_box(slide, Inches(0.5), Inches(3.25), Inches(5.5), Inches(1.3), RGBColor(0xD1, 0xFA, 0xE5), border_color=SUCCESS, border_width=Pt(2))
    tf = be.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, "BACKEND", size=14, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F))
    add_para(tf, "Express.js  ·  TypeScript  ·  MongoDB Atlas", size=10, color=DARK_GRAY)
    add_para(tf, "Firebase Auth  ·  Rate Limiting  ·  43 endpoints  ·  8 collections", size=10, color=DARK_GRAY)

    # AI services boxes
    ai_boxes = [
        ("Ollama\nGemma3 4B", "Local LLM\nInterview Q&A", PRIMARY),
        ("Gemini API\n768-dim vectors", "Cloud\nJob Matching", SUCCESS),
        ("Python Flask\nSTT :5200 / TTS :5100", "Whisper + Edge TTS\nVoice Pipeline", PURPLE),
    ]
    for i, (title, sub, clr) in enumerate(ai_boxes):
        add_flow_box(slide, Inches(0.5 + i * 2.0), Inches(4.75), title, sub, clr, w=Inches(1.8), h=Inches(1.3))

    # Arrows down to AI
    for i in range(3):
        add_arrow(slide, Inches(1.1 + i * 2.0), Inches(4.55), Inches(0.5), Inches(0.2), DARK_GRAY)

    # -- 6 AI Systems table --
    ai_data = [
        ["#", "AI System", "Size", "Where", "What It Does"],
        ["1", "Gemma3 4B (Ollama)", "4.3B params", "Local GPU", "Interview conversation, question gen, feedback"],
        ["2", "Gemini embedding-001", "768-dim", "Google Cloud", "Semantic vectors for job matching"],
        ["3", "Whisper medium", "~500M INT8", "Local CPU", "Speech-to-text (EN/HI/MR + 7 more)"],
        ["4", "Edge TTS (Microsoft)", "Neural voices", "Cloud", "Text-to-speech (3 Indian voices)"],
        ["5", "face-api.js", "892 KB", "Browser", "Face detection + gaze + proctoring"],
        ["6", "GPT-4o (OnDemand)", "Cloud", "Cloud API", "Resume parsing + skill gap explanation"],
    ]
    add_table(slide, Inches(6.4), Inches(1.3), Inches(6.5), Inches(4.5), ai_data,
              col_widths=[Inches(0.4), Inches(1.5), Inches(0.9), Inches(1.0), Inches(2.7)])

    # Key stat
    ks = add_shape_box(slide, Inches(6.4), Inches(6.0), Inches(6.5), Inches(0.5), PRIMARY)
    set_text(ks.text_frame, "  6 AI models  ·  3 languages  ·  <5s voice response  ·  1 unified pipeline", size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_03_resume(prs):
    """SLIDE 3: Feature A — Resume to Job Recommendations (STAR)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Feature A — Resume to Job Recommendations", "STAR: Situation  →  Task  →  Action  →  Result")

    # Situation + Task (left column)
    st_box = add_shape_box(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(1.1), WHITE, border_color=DANGER, border_width=Pt(2))
    tf = st_box.text_frame
    tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.08); tf.word_wrap = True
    set_text(tf, "SITUATION + TASK", size=12, bold=True, color=DANGER)
    add_para(tf, "88% of employers say ATS rejects qualified talent due to keyword mismatch (Harvard, 2021).", size=10, color=DARK_GRAY)
    add_para(tf, "Task: PDF upload → structured profile → semantic job match with explainable scores.", size=10, color=DARK_GRAY)

    # Action: Pipeline flow diagram
    act_label = add_shape_box(slide, Inches(0.4), Inches(2.45), Inches(1.2), Inches(0.35), WARNING)
    set_text(act_label.text_frame, "  ACTION", size=11, bold=True, color=WHITE)

    pipeline_steps = [
        ("PDF.js\nExtract Text", "~0.5s", RGBColor(0x33, 0x99, 0xCC)),
        ("Regex + GPT-4o\nHybrid Parse", "~2.5s", RGBColor(0x66, 0x66, 0xCC)),
        ("Gemini\n768-dim x3", "~1s", SUCCESS),
        ("Atlas Vector\nSearch (HNSW)", "<100ms", PRIMARY),
        ("Cosine\nRerank Top 10", "<10ms", PURPLE),
    ]
    for i, (title, time, clr) in enumerate(pipeline_steps):
        x = Inches(0.4 + i * 2.55)
        add_flow_box(slide, x, Inches(2.9), title, time, clr, w=Inches(2.1), h=Inches(1.0))
        if i < len(pipeline_steps) - 1:
            add_arrow(slide, x + Inches(2.15), Inches(3.25), Inches(0.35), Inches(0.2), clr)

    # Matching formula box
    formula_box = add_shape_box(slide, Inches(0.4), Inches(4.15), Inches(6.4), Inches(2.1), RGBColor(0xF0, 0xF4, 0xFF), border_color=PRIMARY, border_width=Pt(2))
    tf = formula_box.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, "Matching Formula (Research-Backed Weights)", size=13, bold=True, color=PRIMARY)
    add_para(tf, "Match = Skills x 50%  +  Experience x 30%  +  Role Fit x 20%", size=14, bold=True, color=NEAR_BLACK, space_before=Pt(8))
    add_para(tf, "", size=6, color=NEAR_BLACK)
    add_para(tf, "Skills (r=0.54)      → 0.54/1.07 ≈ 50%   — Schmidt & Hunter (1998)", size=10, color=DARK_GRAY)
    add_para(tf, "Experience (r=0.33)  → 0.33/1.07 ≈ 30%   — 85-year meta-analysis", size=10, color=DARK_GRAY)
    add_para(tf, "Role Fit (r=0.20)     → 0.20/1.07 ≈ 20%   — Kristof-Brown (2005)", size=10, color=DARK_GRAY)
    add_para(tf, "Deloitte (2022): Skills-based orgs 107% more likely to place talent", size=10, bold=True, color=PRIMARY, space_before=Pt(6))

    # Result box
    res_box = add_shape_box(slide, Inches(7.1), Inches(4.15), Inches(5.8), Inches(2.1), RGBColor(0xE6, 0xFC, 0xF5), border_color=SUCCESS, border_width=Pt(2))
    tf = res_box.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, "RESULT", size=13, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F))
    add_para(tf, "Resume parsed in ~3 seconds (95% confidence)", size=11, color=NEAR_BLACK, space_before=Pt(8))
    add_para(tf, "Job matching in <100ms (HNSW: O(log N))", size=11, color=NEAR_BLACK)
    add_para(tf, "Each match shows: overall %, skills %, exp %, role fit %", size=11, color=NEAR_BLACK)
    add_para(tf, "Employer side: same algorithm ranks candidates", size=11, color=NEAR_BLACK)
    add_para(tf, "Hidden talent surfaced — 88% that keyword ATS misses", size=11, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F), space_before=Pt(6))

    # Semantic vs Keyword illustration
    vs_box = add_shape_box(slide, Inches(7.1), Inches(1.2), Inches(5.8), Inches(2.7), WHITE, border_color=PURPLE, border_width=Pt(2))
    tf = vs_box.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, "Why Semantic > Keywords", size=13, bold=True, color=PURPLE)
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, "Traditional ATS (keyword):", size=11, bold=True, color=DANGER, space_before=Pt(6))
    add_para(tf, 'Job: "React"  |  Resume: "Built SPAs with React.js, Redux"', size=10, color=DARK_GRAY)
    add_para(tf, "Result: REJECTED (format mismatch)", size=10, bold=True, color=DANGER)
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, "YuvaSetu (semantic embedding):", size=11, bold=True, color=SUCCESS, space_before=Pt(4))
    add_para(tf, "Job vector: [0.23, -0.41, 0.87, ...] (768 dims)", size=10, color=DARK_GRAY)
    add_para(tf, "Resume vector: [0.21, -0.39, 0.85, ...]", size=10, color=DARK_GRAY)
    add_para(tf, "Cosine similarity: 0.94 → 94% MATCH ✓", size=10, bold=True, color=SUCCESS)


def slide_04_skillgap(prs):
    """SLIDE 4: Feature B — Skill Gap Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Feature B — Skill Gap Analysis", "STAR: Situation  →  Task  →  Action  →  Result")

    # Situation
    st = add_shape_box(slide, Inches(0.4), Inches(1.2), Inches(6.2), Inches(0.8), WHITE, border_color=DANGER, border_width=Pt(2))
    tf = st.text_frame; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06); tf.word_wrap = True
    set_text(tf, "SITUATION:", size=11, bold=True, color=DANGER)
    add_para(tf, "Students don't know which skills they lack. Free govt courses (NPTEL: 3,353 / SWAYAM: 11,772) exist but are scattered.", size=10, color=DARK_GRAY)

    # Action: 3-step flow
    act_label = add_shape_box(slide, Inches(0.4), Inches(2.15), Inches(1.0), Inches(0.3), WARNING)
    set_text(act_label.text_frame, " ACTION", size=10, bold=True, color=WHITE)

    # Step 1
    s1 = add_flow_box(slide, Inches(0.4), Inches(2.6), "Step 1: Math Score\n(Cosine Similarity)", "NOT hallucinated\nDeterministic", PRIMARY, w=Inches(3.8), h=Inches(1.1))
    add_arrow(slide, Inches(4.3), Inches(3.0), Inches(0.4), Inches(0.2), PRIMARY)
    # Step 2
    s2 = add_flow_box(slide, Inches(4.8), Inches(2.6), "Step 2: GPT-4o\nExplains the Gap", '"Score is 72%.\nMissing: React, Docker"', PURPLE, w=Inches(3.8), h=Inches(1.1))
    add_arrow(slide, Inches(8.7), Inches(3.0), Inches(0.4), Inches(0.2), PURPLE)
    # Step 3
    s3 = add_flow_box(slide, Inches(9.2), Inches(2.6), "Step 3: Course\nRecommendation", "237 govt courses\n+ Udemy fallback", SUCCESS, w=Inches(3.7), h=Inches(1.1))

    # Key design callout
    kd = add_shape_box(slide, Inches(0.4), Inches(3.9), Inches(12.5), Inches(0.5), RGBColor(0xFF, 0xF3, 0xCD), border_color=WARNING, border_width=Pt(1))
    set_text(kd.text_frame, "  KEY DESIGN: Score is math-computed (cosine similarity), NOT LLM-generated. AI is anchored: \"The score is 72%. Explain WHY. Do NOT recalculate.\"", size=10, bold=True, color=RGBColor(0x85, 0x6D, 0x0E))

    # Course signal scoring table
    signal_data = [
        ["Signal", "Points", "Example"],
        ["Exact skill in course tags", "+50", '"Python" in skills array'],
        ["Full phrase in title", "+40", '"Machine Learning" in title'],
        ["Phrase in description", "+15", '"SQL" in description text'],
        ["Partial word match (>50%)", "+20", '2+ words of "NLP" match'],
        ["Free course bonus", "+2", "Government courses prioritized"],
    ]
    add_table(slide, Inches(0.4), Inches(4.6), Inches(6.2), Inches(2.2), signal_data,
              col_widths=[Inches(2.4), Inches(0.8), Inches(3.0)])

    # Course sources + Result
    src_data = [
        ["Provider", "Courses", "National Scale"],
        ["NPTEL (IITs + IISc)", "119 scraped", "Part of SWAYAM"],
        ["SWAYAM (Govt of India)", "111 scraped", "1.21 Crore (12.1M) users"],
        ["Skill India Digital Hub", "7 curated", "1.5 Crore (15M) candidates"],
        ["Total in DB", "237", "All free, govt-backed"],
    ]
    add_table(slide, Inches(6.9), Inches(4.6), Inches(5.95), Inches(2.2), src_data,
              col_widths=[Inches(2.2), Inches(1.3), Inches(2.45)], header_color=SUCCESS)


def slide_05_mocktest(prs):
    """SLIDE 5: Feature C — Mock Test"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Feature C — Mock Test (FAANG Questions)", "STAR: Situation  →  Task  →  Action  →  Result")

    # Situation + Task
    st = add_shape_box(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.7), WHITE, border_color=DANGER, border_width=Pt(2))
    tf = st.text_frame; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06); tf.word_wrap = True
    set_text(tf, "SITUATION: Students lack structured practice material. Resources are scattered with no tracking.", size=11, bold=False, color=DARK_GRAY)
    add_para(tf, "TASK: Curated question bank + timed tests + progress tracking across attempts.", size=11, bold=False, color=DARK_GRAY)

    # Action: 3-phase flow
    phases = [
        ("Question Bank", "300+ MCQs\nFAANG-style\nSystem Design\nAlgorithms", PRIMARY),
        ("Timed Quiz", "Structured session\nImmediate scoring\nTopic breakdown", WARNING),
        ("Performance", "Score tracking\nTrend analysis\n↑ / → / ↓", SUCCESS),
    ]
    for i, (title, desc, clr) in enumerate(phases):
        x = Inches(0.5 + i * 4.3)
        box = add_shape_box(slide, x, Inches(2.3), Inches(3.8), Inches(2.0), WHITE, border_color=clr, border_width=Pt(2))
        tf = box.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
        set_text(tf, title, size=16, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, desc, size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
        if i < 2:
            add_arrow(slide, x + Inches(3.85), Inches(3.1), Inches(0.4), Inches(0.2), clr)

    # Bar chart: topics
    categories = ["Sys Design", "Algorithms", "Data Struct", "OOP", "Databases"]
    values = (85, 70, 65, 50, 45)
    add_bar_chart(slide, Inches(0.5), Inches(4.6), Inches(5.5), Inches(2.6), categories, values, "Questions by Topic", PRIMARY)

    # Result box
    res = add_shape_box(slide, Inches(6.5), Inches(4.6), Inches(6.4), Inches(2.6), RGBColor(0xE6, 0xFC, 0xF5), border_color=SUCCESS, border_width=Pt(2))
    tf = res.text_frame; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15); tf.word_wrap = True
    set_text(tf, "RESULT", size=14, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F))
    add_para(tf, "300+ curated FAANG questions with answers", size=12, color=NEAR_BLACK, space_before=Pt(10))
    add_para(tf, "Timed sessions with immediate scoring", size=12, color=NEAR_BLACK)
    add_para(tf, "Topic-level weakness identification", size=12, color=NEAR_BLACK)
    add_para(tf, "Trend analysis: improving / stable / declining", size=12, color=NEAR_BLACK)
    add_para(tf, "Complements voice AI interview practice", size=12, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F), space_before=Pt(8))


def slide_06_interview(prs):
    """SLIDE 6: Feature D — AI Mock Interview + Proctoring"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Feature D — AI Mock Interview + Proctoring", "STAR: Voice-based AI interviewer in Hindi/Marathi/English with live proctoring")

    # Situation
    st = add_shape_box(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.55), WHITE, border_color=DANGER, border_width=Pt(2))
    tf = st.text_frame; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04); tf.word_wrap = True
    set_text(tf, "SITUATION: Coaching costs Rs 5K-50K. No voice AI in Indian languages. AI practice → 30% communication improvement (Forbes, 2023)", size=11, color=DARK_GRAY)

    # Voice pipeline flow
    pipeline_label = add_shape_box(slide, Inches(0.4), Inches(1.9), Inches(1.6), Inches(0.3), PRIMARY)
    set_text(pipeline_label.text_frame, " Voice Pipeline", size=10, bold=True, color=WHITE)

    steps = [
        ("You Speak\nPCM → WAV", "~100ms", RGBColor(0x64, 0x74, 0x8B)),
        ("Whisper STT\n+ VAD + filter", "~2s", RGBColor(0x33, 0x99, 0xCC)),
        ("Gemma3 LLM\n1-2 sentences", "~2s", PRIMARY),
        ("Edge TTS\nparallel synth", "~0.3s", PURPLE),
        ("You Hear\nAI responds", "", SUCCESS),
    ]
    for i, (title, time, clr) in enumerate(steps):
        x = Inches(0.4 + i * 2.55)
        add_flow_box(slide, x, Inches(2.3), title, time, clr, w=Inches(2.2), h=Inches(0.9))
        if i < len(steps) - 1:
            add_arrow(slide, x + Inches(2.25), Inches(2.6), Inches(0.25), Inches(0.18), clr)

    # Total latency callout
    lat = add_shape_box(slide, Inches(0.4), Inches(3.3), Inches(12.5), Inches(0.35), RGBColor(0xDB, 0xEA, 0xFE))
    set_text(lat.text_frame, "  Total: ~4-5 seconds from your voice to AI's voice reply  |  10 personalized questions  |  Conversational (not scripted)", size=10, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # Proctoring box (left)
    proc = add_shape_box(slide, Inches(0.4), Inches(3.85), Inches(5.8), Inches(3.2), WHITE, border_color=DANGER, border_width=Pt(2))
    tf = proc.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, "Proctoring (100% client-side — no video sent to server)", size=12, bold=True, color=DANGER)
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, "face-api.js (892 KB)  ·  Detects every 500ms", size=10, bold=True, color=DARK_GRAY, space_before=Pt(4))
    add_para(tf, "", size=4, color=DARK_GRAY)
    violations = [
        "No face detected       →  -5 pts    (2s grace period)",
        "Multiple faces            →  -15 pts  (instant)",
        "Looking away             →  -3 pts    (2s grace, 4-method gaze)",
        "Tab switch                  →  -10 pts  (instant, 100% accurate)",
    ]
    for v in violations:
        add_para(tf, v, size=10, color=NEAR_BLACK)
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, "3 Strikes = AUTO-TERMINATE", size=12, bold=True, color=DANGER, space_before=Pt(4))
    add_para(tf, "Integrity Score = 100 - total_deductions", size=10, color=DARK_GRAY)
    add_para(tf, "Gaze: head pose + face symmetry + eye width + iris analysis", size=9, color=SECONDARY, space_before=Pt(4))

    # Feedback box (right)
    fb = add_shape_box(slide, Inches(6.5), Inches(3.85), Inches(6.4), Inches(3.2), WHITE, border_color=SUCCESS, border_width=Pt(2))
    tf = fb.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, "5-Category Scored Feedback", size=12, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F))
    add_para(tf, "", size=4, color=DARK_GRAY)
    cats = [
        ("Communication", "78", "Clarity, structure, language"),
        ("Technical", "72", "Accuracy, depth, decisions"),
        ("Problem Solving", "65", "Structured thinking, examples"),
        ("Cultural Fit", "80", "Teamwork, enthusiasm"),
        ("Confidence", "70", "Assertiveness, handling pressure"),
    ]
    for name, score, desc in cats:
        bar_len = int(int(score) * 0.20)
        bar = "█" * bar_len + "░" * (20 - bar_len)
        add_para(tf, f"{name:20s} {bar} {score}/100", size=9, color=NEAR_BLACK, font_name="Consolas")
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, "Total: 73/100 (Good)", size=12, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F), space_before=Pt(4))
    add_para(tf, "+ 3 Strengths  ·  + 3 Improvements  ·  + Final Assessment", size=10, color=DARK_GRAY)
    add_para(tf, "Available in: English  |  Hindi  |  Marathi", size=10, bold=True, color=PRIMARY, space_before=Pt(4))


def slide_07_cost(prs):
    """SLIDE 7: Production Cost Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Production Cost Analysis — GCP Mumbai Region", "GKE Standard  ·  Spot VMs  ·  KEDA Scale-to-Zero  ·  Pay-as-you-go")

    # Architecture diagram
    arch_label = add_shape_box(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.35), PRIMARY)
    set_text(arch_label.text_frame, "  2-Pool GKE Standard Architecture", size=11, bold=True, color=WHITE)

    # Pool 1
    p1 = add_shape_box(slide, Inches(0.4), Inches(1.65), Inches(6.0), Inches(1.3), RGBColor(0xD1, 0xFA, 0xE5), border_color=SUCCESS, border_width=Pt(2))
    tf = p1.text_frame; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06); tf.word_wrap = True
    set_text(tf, 'Pool 1: "Always On" — e2-standard-2 (2 vCPU, 8GB)', size=11, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F))
    add_para(tf, "Express Backend  ·  React Frontend  ·  MongoDB  ·  Qdrant Vector DB", size=10, color=DARK_GRAY)
    add_para(tf, "Spot VM:  Rs 1,800/mo ($22)  —  runs 24/7", size=10, bold=True, color=NEAR_BLACK)

    # Pool 2
    p2 = add_shape_box(slide, Inches(0.4), Inches(3.1), Inches(6.0), Inches(1.3), RGBColor(0xDB, 0xEA, 0xFE), border_color=PRIMARY, border_width=Pt(2))
    tf = p2.text_frame; tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06); tf.word_wrap = True
    set_text(tf, 'Pool 2: "Scale to Zero" — g2-standard-4 + NVIDIA L4 GPU', size=11, bold=True, color=PRIMARY)
    add_para(tf, "Ollama Gemma3  ·  Whisper STT  ·  Edge TTS", size=10, color=DARK_GRAY)
    add_para(tf, "Spot GPU:  Rs 12/hr  —  0 nodes when idle = Rs 0 cost  (KEDA autoscale)", size=10, bold=True, color=NEAR_BLACK)

    # Per user per day
    user_data = [
        ["Feature", "Service", "Cost (Rs)"],
        ["Mock Interview (20 min)", "L4 GPU Spot", "2.50 - 4.00"],
        ["Skill Gap Analysis", "L4 GPU Spot", "0.25 - 0.40"],
        ["Job Matching", "CPU + MongoDB", "0.10 - 0.20"],
        ["Vector Embeddings", "CPU only", "0.05 - 0.15"],
        ["Storage & DB", "MongoDB + GCS", "0.60 - 1.20"],
        ["DAILY TOTAL", "", "Rs 3.50 - 6.00"],
    ]
    add_table(slide, Inches(6.8), Inches(1.2), Inches(6.1), Inches(2.8), user_data,
              col_widths=[Inches(2.4), Inches(1.6), Inches(2.1)])

    # Scaling comparison chart
    scale_data = [
        ["Scale", "Autopilot", "Standard", "Savings"],
        ["5 users", "Rs 4,500", "Rs 4,200", "~7%"],
        ["500 users", "Rs 18,000", "Rs 7,500", "58%"],
        ["5,000 users", "Rs 85,000", "Rs 25,000", "70%"],
    ]
    add_table(slide, Inches(0.4), Inches(4.6), Inches(6.0), Inches(1.5), scale_data,
              col_widths=[Inches(1.3), Inches(1.5), Inches(1.5), Inches(1.7)], header_color=PURPLE)

    # The killer comparison
    comp = add_shape_box(slide, Inches(6.8), Inches(4.3), Inches(6.1), Inches(1.8), WHITE, border_color=SUCCESS, border_width=Pt(3))
    tf = comp.text_frame; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15); tf.word_wrap = True
    set_text(tf, "The Comparison That Matters", size=14, bold=True, color=RGBColor(0x0B, 0x6E, 0x4F))
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, "YuvaSetu:    Rs 5/day per user", size=16, bold=True, color=SUCCESS, space_before=Pt(6))
    add_para(tf, "Traditional: Rs 5,000-50,000 per session", size=14, bold=True, color=DANGER, space_before=Pt(6))
    add_para(tf, "> 1000x cheaper. Same features. In Hindi.", size=12, bold=True, color=PRIMARY, space_before=Pt(8))

    # Fixed baseline note
    note = add_shape_box(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.5), RGBColor(0xF8, 0xF9, 0xFA))
    set_text(note.text_frame, "  Fixed baseline (zero users): Rs 3,750/mo ($45) — LB + system node + disk. GKE fee covered by $74.40 GCP credit.", size=9, color=SECONDARY)


def slide_08_different(prs):
    """SLIDE 8: Why It's Different"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Why YuvaSetu Is Different", "No existing platform combines all these features")

    # Comparison matrix
    comp_data = [
        ["Feature", "YuvaSetu", "Unstop", "Pramp", "InterviewBit", "Naukri", "LinkedIn"],
        ["Voice AI Interview", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["Hindi/Marathi Voice", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["AI Proctoring", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["Semantic Job Match", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["Skill Gap → Courses", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["Free Govt Courses", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["Local LLM (Rs 0)", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["5-Cat Scored Feedback", "✓", "✗", "✗", "✗", "✗", "✗"],
        ["Hiring Challenges", "✗", "✓", "✗", "✗", "✗", "✗"],
    ]
    add_table(slide, Inches(0.4), Inches(1.2), Inches(7.4), Inches(3.8), comp_data,
              col_widths=[Inches(1.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(1.0)])

    # Differentiator cards (right side)
    diffs = [
        ("Local LLM", "Data NEVER leaves machine.\nZero API cost. Privacy-first.", PRIMARY),
        ("Semantic Matching", "Surfaces 88% of talent\nthat keyword ATS misses.", PURPLE),
        ("Skills-First Weights", "Skills-based orgs 107% more\nlikely to place talent (Deloitte).", SUCCESS),
        ("Govt Course Integration", "237 free NPTEL/SWAYAM courses.\nNot Udemy upsells.", WARNING),
        ("Client-Side Proctoring", "No video to server. 892KB\nmodels run in browser.", DANGER),
        ("End-to-End Pipeline", "Resume → Match → Gap →\nLearn → Interview → Feedback.", RGBColor(0x64, 0x74, 0x8B)),
    ]
    for i, (title, desc, clr) in enumerate(diffs):
        row = i // 2
        col = i % 2
        x = Inches(7.7 + col * 2.9)
        y = Inches(1.2 + row * 1.7)
        card = add_shape_box(slide, x, y, Inches(2.7), Inches(1.5), WHITE, border_color=clr, border_width=Pt(2))
        tf = card.text_frame; tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.08); tf.word_wrap = True
        set_text(tf, title, size=11, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, desc, size=9, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))

    # Hidden talent callout
    ht = add_shape_box(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.8), RGBColor(0xF5, 0xF0, 0xFF), border_color=PURPLE, border_width=Pt(2))
    tf = ht.text_frame; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.12); tf.word_wrap = True
    set_text(tf, 'The "Hidden Talent" Problem We Solve', size=14, bold=True, color=PURPLE)
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, 'Traditional ATS:  Job requires "React" → Resume says "Built SPAs with React.js, Redux, TypeScript" → REJECTED (keyword format mismatch)', size=11, color=DANGER, space_before=Pt(6))
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, 'YuvaSetu Semantic:  Job vector [768-dim] vs Resume vector [768-dim] → Cosine similarity: 0.94 → 94% MATCH (understands meaning, not just words)', size=11, color=SUCCESS, space_before=Pt(4))
    add_para(tf, "", size=4, color=DARK_GRAY)
    add_para(tf, '88% of employers believe their ATS rejects qualified candidates — Harvard Business School (2021)', size=11, bold=True, color=PURPLE, space_before=Pt(4))


def slide_09_research(prs):
    """SLIDE 9: Research & References"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Research & Evidence Base", "Peer-reviewed papers  ·  Government data  ·  Industry reports")

    # Employment data
    emp_header = add_shape_box(slide, Inches(0.4), Inches(1.15), Inches(6.2), Inches(0.3), DANGER)
    set_text(emp_header.text_frame, "  A. Employment Crisis Data", size=10, bold=True, color=WHITE)
    emp_data = [
        ["Statistic", "Value", "Source"],
        ["Engg graduates/year", "1.5M (10% get jobs)", "Business Standard, 2024"],
        ["Unemployed after degree", "83%", "Unstop Talent Report 2025"],
        ["Grads under 25 jobless", "40%", "Azim Premji Univ, 2026"],
        ["Employers: skill mismatch", "82%", "India Skills Report 2025"],
        ["Graduate employability", "42.6%", "Mercer-Mettl GSI 2025"],
    ]
    add_table(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(2.2), emp_data,
              col_widths=[Inches(2.0), Inches(1.6), Inches(2.6)], header_color=DANGER)

    # Algorithm research
    algo_header = add_shape_box(slide, Inches(6.9), Inches(1.15), Inches(6.0), Inches(0.3), PRIMARY)
    set_text(algo_header.text_frame, "  B. Algorithm & Matching Research", size=10, bold=True, color=WHITE)
    algo_data = [
        ["Paper", "Key Finding", "How We Use It"],
        ["Schmidt & Hunter (1998)", "Skills tests: r=0.54 highest", "Skills weighted 50%"],
        ["Deloitte (2022)", "Skills-based: 107% better hiring", "Skills as primary vector"],
        ["Holistic Triangle (2024)", "3 dims: skill + exp + fit", "3-vector approach"],
        ["Harvard (2021)", "88% ATS rejects good talent", "Semantic matching"],
        ["arXiv (2021)", "Multi-vector > keywords", "Cosine reranking"],
    ]
    add_table(slide, Inches(6.9), Inches(1.5), Inches(6.0), Inches(2.2), algo_data,
              col_widths=[Inches(1.8), Inches(2.0), Inches(2.2)], header_color=PRIMARY)

    # AI Interview research
    ai_header = add_shape_box(slide, Inches(0.4), Inches(3.9), Inches(6.2), Inches(0.3), PURPLE)
    set_text(ai_header.text_frame, "  C. AI Interview Research", size=10, bold=True, color=WHITE)
    ai_data = [
        ["Finding", "Statistic", "Source"],
        ["AI practice improves comm.", "30% improvement", "Forbes HR Council, 2023"],
        ["Low-pressure AI simulation", "Reduces interview anxiety", "Predictive Validity, 2022"],
        ["Embedding cosine sim", "r > 0.80 vs human judgment", "Sentence-BERT, 2019"],
    ]
    add_table(slide, Inches(0.4), Inches(4.25), Inches(6.2), Inches(1.5), ai_data,
              col_widths=[Inches(2.2), Inches(1.8), Inches(2.2)], header_color=PURPLE)

    # Govt platform data
    govt_header = add_shape_box(slide, Inches(6.9), Inches(3.9), Inches(6.0), Inches(0.3), SUCCESS)
    set_text(govt_header.text_frame, "  D. Government Platform Data", size=10, bold=True, color=WHITE)
    govt_data = [
        ["Platform", "Users", "Courses"],
        ["SWAYAM", "1.21 Crore (12.1M)", "11,772"],
        ["Skill India (SIDH)", "1.5 Crore (15M)", "752"],
        ["NPTEL (IITs)", "Part of SWAYAM", "3,353"],
    ]
    add_table(slide, Inches(6.9), Inches(4.25), Inches(6.0), Inches(1.5), govt_data,
              col_widths=[Inches(1.8), Inches(2.0), Inches(2.2)], header_color=SUCCESS)

    # Weight derivation box
    wd = add_shape_box(slide, Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.1), RGBColor(0xF0, 0xF4, 0xFF), border_color=PRIMARY, border_width=Pt(2))
    tf = wd.text_frame; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.08); tf.word_wrap = True
    set_text(tf, "Weight Derivation (50/30/20) — Proportional to Meta-Analytic Validity Coefficients", size=11, bold=True, color=PRIMARY)
    add_para(tf, "Skills: r ≈ 0.54 → 0.54/1.07 = 50%   |   Experience: r ≈ 0.33 → 0.33/1.07 = 30%   |   Role Fit: r ≈ 0.20 → 0.20/1.07 = 20%", size=10, color=NEAR_BLACK, space_before=Pt(6))
    add_para(tf, "Sources: Schmidt & Hunter 85-year meta-analysis (1998)  ·  Kristof-Brown et al. (2005)  ·  Quinones, Ford & Teachout (1995)", size=9, color=SECONDARY, space_before=Pt(4))


def slide_10_impact(prs):
    """SLIDE 10: Impact + Demo Script"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_WHITE)
    add_title_bar(slide, "Impact + Live Demo", "YuvaSetu — AI-Powered Career Platform for Bharat")

    # Stats banner
    stats = [
        ("28,680", "Lines of\nCode", PRIMARY),
        ("6", "AI Systems\nIntegrated", PURPLE),
        ("43", "API\nEndpoints", SUCCESS),
        ("3", "Languages\nSupported", WARNING),
        ("237", "Govt Courses\nScraped", RGBColor(0x0B, 0x6E, 0x4F)),
        ("Rs 0", "Interview\nCost", DANGER),
    ]
    for i, (num, lbl, clr) in enumerate(stats):
        add_stat_card(slide, Inches(0.4 + i * 2.15), Inches(1.2), num, lbl, clr, width=Inches(1.95), height=Inches(1.2))

    # Before vs After
    before_after = [
        ["Current Reality", "With YuvaSetu"],
        ["Coaching: Rs 5,000-50,000/session", "Rs 0 (local AI)"],
        ["English-only platforms", "Hindi + Marathi + English"],
        ["No interview feedback", "5-category scored analysis"],
        ["Blind job applications", "AI-matched with % breakdown"],
        ["Google search for courses", "237 verified govt course links"],
        ["No cheating detection", "AI proctoring (face + gaze + tab)"],
        ["88% talent rejected by ATS", "Semantic matching finds them"],
    ]
    add_table(slide, Inches(0.4), Inches(2.65), Inches(6.2), Inches(3.2), before_after,
              col_widths=[Inches(3.1), Inches(3.1)], header_color=RGBColor(0x34, 0x34, 0x34))

    # Demo script
    demo_header = add_shape_box(slide, Inches(6.9), Inches(2.65), Inches(6.0), Inches(0.35), PRIMARY)
    set_text(demo_header.text_frame, "  Live Demo Script (5 minutes)", size=11, bold=True, color=WHITE)

    demo_data = [
        ["Time", "Action", "What Judges See"],
        ["0:00", "Upload resume", "AI parses in 3s, 95% confidence"],
        ["1:00", "Job matches", "Top jobs with % breakdown"],
        ["2:00", "Skill gap", "Missing skills + NPTEL courses"],
        ["3:00", "Hindi interview", "AI voice conversation"],
        ["3:30", "Trigger proctor", "Look away → warning"],
        ["4:00", "End interview", "5-cat feedback in Hindi"],
        ["4:30", "Employer view", "AI-ranked candidates"],
        ["5:00", "Close", "Impact statement"],
    ]
    add_table(slide, Inches(6.9), Inches(3.05), Inches(6.0), Inches(2.8), demo_data,
              col_widths=[Inches(0.7), Inches(1.6), Inches(3.7)])

    # Closing one-liner
    close = add_shape_box(slide, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.8), PRIMARY)
    tf = close.text_frame; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.1); tf.word_wrap = True
    set_text(tf, '"No other platform lets you talk to an AI interviewer in Hindi while it monitors you', size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_para(tf, 'with face detection and recommends free government courses — all for Rs 5/day."', size=13, bold=True, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER, space_before=Pt(2))


# ═══════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_problem(prs)
    slide_02_tech(prs)
    slide_03_resume(prs)
    slide_04_skillgap(prs)
    slide_05_mocktest(prs)
    slide_06_interview(prs)
    slide_07_cost(prs)
    slide_08_different(prs)
    slide_09_research(prs)
    slide_10_impact(prs)

    output_path = r"c:\Users\Hp\OneDrive\Desktop\Algozenith\CareerBridge\Hack_DTU_Main\YuvaSetu_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
